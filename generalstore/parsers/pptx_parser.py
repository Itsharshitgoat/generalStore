"""
PPTX parser using python-pptx.

Each slide produces one DocumentChunk.  The slide title (if any) is
placed first in the chunk text, followed by text from all other shapes
that have a text_frame.
"""

import logging
from pathlib import Path
from typing import List, Optional

from pptx import Presentation

from generalstore.config import get_settings
from generalstore.parsers.base import BaseParser, DocumentChunk

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """Parser for .pptx presentation files."""

    supported_extensions: List[str] = [".pptx"]

    def parse(self, filepath: Path) -> List[DocumentChunk]:
        """
        Parse a PPTX file and return a list of DocumentChunks.

        Each slide is emitted as a single chunk.  The slide title
        (from ``slide.shapes.title``) is prepended to the combined
        text of every shape that has a text_frame.

        Args:
            filepath: Absolute path to the .pptx file.

        Returns:
            List of DocumentChunk objects with slide-level metadata.

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If a .ppt file is supplied (unsupported format).
            RuntimeError: If PPTX parsing fails.
        """
        filepath = Path(filepath).resolve()
        if not filepath.exists():
            raise FileNotFoundError(f"PPTX file not found: {filepath}")

        if filepath.suffix.lower() == ".ppt":
            logger.warning(
                "Legacy .ppt format is not supported: '%s'. "
                "Please convert to .pptx.",
                filepath.name,
            )
            raise ValueError(
                f"Unsupported format '.ppt' for file '{filepath.name}'. "
                "Only .pptx is supported."
            )

        settings = get_settings()
        subject = self._derive_subject(filepath)
        chunks: List[DocumentChunk] = []
        chunk_index = 0

        try:
            prs = Presentation(str(filepath))
        except Exception as exc:
            raise RuntimeError(
                f"Failed to open PPTX '{filepath.name}': {exc}"
            ) from exc

        for slide_idx, slide in enumerate(prs.slides):
            slide_number = slide_idx + 1  # 1-indexed

            try:
                title = self._extract_title(slide)
                body_parts = self._extract_body_texts(slide)

                # Build combined text: title first, then body
                parts: List[str] = []
                if title:
                    parts.append(title)
                parts.extend(body_parts)

                combined = "\n".join(parts)
                cleaned = self._clean_text(combined)

                if len(cleaned) < settings.min_chunk_chars:
                    continue

                # Respect max_chunk_chars
                segments = self._split_text(cleaned, settings.max_chunk_chars)
                for seg in segments:
                    if len(seg) < settings.min_chunk_chars:
                        continue
                    chunks.append(
                        DocumentChunk(
                            text=seg,
                            source_file=str(filepath),
                            file_type=".pptx",
                            chunk_index=chunk_index,
                            slide_number=slide_number,
                            heading=title,
                            subject=subject,
                        )
                    )
                    chunk_index += 1

            except Exception as exc:
                logger.warning(
                    "Failed to parse slide %d of '%s': %s",
                    slide_number,
                    filepath.name,
                    exc,
                )

        logger.info(
            "Parsed PPTX '%s': %d chunks from %d slides",
            filepath.name,
            len(chunks),
            len(prs.slides),
        )
        return chunks

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _extract_title(slide) -> Optional[str]:
        """Return the slide's title text, or None."""
        if slide.shapes.title is not None:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                return title_text
        return None

    @staticmethod
    def _extract_body_texts(slide) -> List[str]:
        """Collect text from all shapes with a text_frame (excluding the title)."""
        title_shape = slide.shapes.title
        texts: List[str] = []

        for shape in slide.shapes:
            # Skip the title shape – already handled
            if shape is title_shape:
                continue
            if shape.has_text_frame:
                frame_text = shape.text_frame.text.strip()
                if frame_text:
                    texts.append(frame_text)
        return texts

    @staticmethod
    def _split_text(text: str, max_chars: int) -> List[str]:
        """Split text into segments of at most *max_chars*, breaking at boundaries."""
        if len(text) <= max_chars:
            return [text]

        parts: List[str] = []
        while text:
            if len(text) <= max_chars:
                parts.append(text)
                break

            split_pos = text.rfind(". ", 0, max_chars)
            if split_pos == -1 or split_pos < max_chars // 2:
                split_pos = text.rfind(" ", 0, max_chars)
            if split_pos == -1:
                split_pos = max_chars

            parts.append(text[: split_pos + 1].strip())
            text = text[split_pos + 1 :].strip()

        return parts
